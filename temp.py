import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import MSO_UNDERLINE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE

# Code RGB du logo : R= 50 G = 171 B = 104


#partie pour mettre la photo dans le coin 
"""
img_path1 = 'picsellia.png'
left = Inches(9.1)
top = Inches(0.15)
width = Inches (0.7)
height = Inches (0.9)
pic = slide. shapes. add_picture(img_path1, left, top, width, height)
"""

def diapo(nom):

#slide 1 

    img_path = 'picsellia.png'
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    date = datetime.datetime.now().strftime("%y-%m-%d")
    img_path1 = 'picsellia.png'
    
    
    title.text = "Rapport d'avancée"
    subtitle.text = "Picsell.ia - 15/07/2019"
    
    left = Inches(4)
    top = Inches(0.2)
    width = Inches (2)
    height = Inches (2.6)
    pic = slide. shapes. add_picture(img_path, left, top, width, height)
    
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    
    #slide2
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    
    #titre
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Contextualisation"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    #j'ai créé plusieurs boites de textes car comme je le disais j'ai pas trouvé comment
    #faire en sorte de retourner à la ligne sans casser la typologie
    #j'ai fait un petit bout de code (à côté) qui me permettait de savoir si 
    #le nombre de caractères était trop grand et ferait dépasser de la ligne et
    #qui disait vers où s'arrêter; ca m'a été un peu utile
    
    #j'ai pris 0.2 pour un retour à la ligne et 0.5 pour un saut de ligne 
    
    left = Inches (0.5)
    top = Inches(1.4)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Rendez-vous pris en mars 2019, nous vous avions présenté notre projet de start-up :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    
    
    left = Inches (0.5)
    top = Inches(1.9)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Nous sommes spécialisés en annotation semi automatique de base de données d'images,"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(2.1)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "de vidéos, de textes et de sons nécessaires au déploiement de vos solutions IA."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(2.6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Nous avions besoin de données pour tester notre service, vous nous avez donc envoyé votre dataset"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    
    left = Inches (0.5)
    top = Inches(2.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "contenant des déchets PET ou non PET / Alimentaire non Alimentaire / Bouteille Barquette."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(3.3)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Ce dataset s'inscrivant dans votre volonté de développer des solutions de tri via intelligence"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(3.5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "artificielle, la labellisation de données n'étant pas votre principal problème car vous avez"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(3.7)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "la main sur vos données, vous nous aviez proposé de vous présenter nos résultats afin que"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(3.9)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "vous puissiez déterminer si notre service répond ou non à vos besoins."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(4.4)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Nous avons toutefois tenu à faire de nombreuses analyses de vos données, afin de vous donner"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(4.6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "le plus d'informations possibles et vous présenter les pistes possibles de développement de votre projet."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    left = Inches (0.5)
    top = Inches(5.1)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Dans un second temps nous allons vous présenter notre service et ce qu'il peut vous apporter."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    
    
    #slide 3 
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Les prérequis de votre projet"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    
    #ici je voulais créer un objet et le mettre ensuite en vert mais des deux facons si desosus
    #j'ai pas réussi parceque les connecteurs et les formes marchent pas directement avec ca
    #du coup à court d'idées et j'ai mis une photo de barre verte
    
    
    #line1=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5), Inches(1.7), Inches(5), Inches(6.7))
    
    #line1=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.7), Inches(5), Inches(6.7))
    #line1.color.rgb = RGBColor(50, 171, 104)
    img_path1 = 'ligne.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1,Inches(5), Inches(1.7), Inches(0.10), Inches(5))
    
    left = Inches (0.5)
    top = Inches(1.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Pertinence de vos données"
    p.font.size = Pt(26)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(0,0,0)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    
    left = Inches (0.5)
    top = Inches(2.1)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "pour votre use-case :"
    p.font.size = Pt(26)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(0,0,0)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    
    left = Inches (0.4)
    top = Inches(3.5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = " - Dataset ayant la même proportion"
    p.font.size = Pt(20)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.4)
    top = Inches(3.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "d'images dans chaque classe."
    p.font.size = Pt(20)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.4)
    top = Inches(4.5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "- Qualité suffisante des images pour"
    p.font.size = Pt(20)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.4)
    top = Inches(4.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "pour obtenir un process de qualité."
    p.font.size = Pt(20)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    
    left = Inches (5.1)
    top = Inches(1.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Faisabilité de votre use-case :"
    p.font.size = Pt(26)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(0,0,0)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    
    
    left = Inches (5.1)
    top = Inches(3.5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "- Bonnes performances de prédiction"
    p.font.size = Pt(20)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50,50,50)
    
    left = Inches (5.1)
    top = Inches(3.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "si l'on a un dataset labellisé."
    p.font.size = Pt(20)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50,50,50)
    
    #slide 4 
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Prérequis 1 : Proportion du datasheet"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    img_path2 = 'C:\codePyth\'',nom,'image_1.png'
    left = Inches(3)
    top = Inches(1.7)
    width = Inches (4)
    height = Inches (3)
    pic = slide. shapes. add_picture(img_path2, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    
    p = tf.add_paragraph()
    p.text = "TTTTTTTTTTTTT :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    
    
    
    #slide5 
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Prérequis 1 : Proportion du datasheet"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    img_path2 = 'C:\codePyth\'',nom,'image_2.png'
    left = Inches(3)
    top = Inches(1.7)
    width = Inches (4)
    height = Inches (3)
    pic = slide. shapes. add_picture(img_path2, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    
    p = tf.add_paragraph()
    p.text = "TTTTTTTTTTTTT :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    
    #Slide 6 
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Prérequis 1 : Proportion du datasheet"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    img_path2 = 'C:\codePyth\'',nom,'image_3.png'
    left = Inches(3)
    top = Inches(1.7)
    width = Inches (4)
    height = Inches (3)
    pic = slide. shapes. add_picture(img_path2, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    
    p = tf.add_paragraph()
    p.text = "TTTTTTTTTTTTT :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    
    
    
    #Slide7
    
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Prérequis 2 : Qualité des images"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    img_path2 ='C:\codePyth\'',nom,'image_4.png'
    left = Inches(2.2)
    top = Inches(1.8)
    width = Inches (2)
    height = Inches (2)
    pic = slide. shapes. add_picture(img_path2, left, top, width, height)
    
    img_path3 = 'C:\codePyth\'',nom,'image_5.png'
    left = Inches(6)
    top = Inches(1.8)
    width = Inches (2)
    height = Inches (2)
    pic = slide. shapes. add_picture(img_path3, left, top, width, height)
    
    
    left = Inches (0.5)
    top = Inches(5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "TTTTTTTTTTTTT :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "PPPPPPPPPPPPPPPP :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.level = 1
    
    
    #Slide8
    
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Prérequis 2 : Qualité des images"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    img_path2 = 'C:\codePyth\'',nom,'image_6.png'
    left = Inches(2.2)
    top = Inches(1.8)
    width = Inches (2)
    height = Inches (2)
    pic = slide. shapes. add_picture(img_path2, left, top, width, height)
    
    img_path3 = 'C:\codePyth\'',nom,'image_7.png'
    left = Inches(6)
    top = Inches(1.8)
    width = Inches (2)
    height = Inches (2)
    pic = slide. shapes. add_picture(img_path3, left, top, width, height)
    
    
    left = Inches (0.5)
    top = Inches(5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "TTTTTTTTTTTTT :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "PPPPPPPPPPPPPPPP :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.level = 1
    
    
    #Slide9
    
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Prérequis 2 : Qualité des images"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    img_path2 = 'C:\codePyth\'',nom,'image_8.png'
    left = Inches(2.2)
    top = Inches(1.8)
    width = Inches (2)
    height = Inches (2)
    pic = slide. shapes. add_picture(img_path2, left, top, width, height)
    
    img_path3 = 'C:\codePyth\'',nom,'image_9.png'
    left = Inches(6)
    top = Inches(1.8)
    width = Inches (2)
    height = Inches (2)
    pic = slide. shapes. add_picture(img_path3, left, top, width, height)
    
    
    left = Inches (0.5)
    top = Inches(5)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "TTTTTTTTTTTTT :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "PPPPPPPPPPPPPPPP :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.level = 1
    
    
    #Slide10
    
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Prérequis 3 : Faisabilité"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    p = tf.add_paragraph()
    p.text = "TTTTTTTTTTTTTTTTTTT :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(127, 127, 127)
    p.font.underline = MSO_UNDERLINE.SINGLE_LINE
    p.level = 1
    
    
    #La slide qui suit c'était un peu plus chiant puisque comme je créais les boites de textes
    #il fallait que je fasse comme du "collage" pour que les différentes couleurs de texte
    #se suivent bien (après c'est censé être une slide qui bouge pas donc ca pause plus trop de pb)
    
    #Slide11
    
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    
    #partie pour mettre la photo dans le coin 
    img_path1 = 'picsellia.png'
    left = Inches(9.1)
    top = Inches(0.15)
    width = Inches (0.7)
    height = Inches (0.9)
    pic = slide. shapes. add_picture(img_path1, left, top, width, height)
    
    left = Inches (0.5)
    top = Inches(0.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Notre service"
    p.font.size = Pt(40)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (0.5)
    top = Inches(1.4)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Nous fournissons un service d'"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (3.12)
    top = Inches(1.4)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "annotation d'image et de conseil en intelligence artificielle."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (0.5)
    top = Inches(1.9)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "En couplant intelligemment l'"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (3)
    top = Inches(1.9)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "action humaine de labellisation et nos algorithmes"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (7.33)
    top = Inches(1.9)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = ", nous optimisons"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(2.1)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "la préparation de votre base de donnée et vous donnons une inside pertinente sur les algorithmes"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(2.3)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "à utiliser pour votre use-case."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(2.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Pour vous donner un ordre d'idée nous, nous labellisons : "
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (5.45)
    top = Inches(2.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "1000 images en 25 minutes."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (0.5)
    top = Inches(3.2)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Ainsi pour créer une database susceptible d'obtenir des bonnes performances pour discriminer"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(3.4)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Bouteille/Barquette il vous faut un dataset d'"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (4.33)
    top = Inches(3.4)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "images différentes"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (5.95)
    top = Inches(3.4)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = """(Sans "data augmentation" -> cf rotation """
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(3.6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "de vos objets lors de la prise de photo) d'environ"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (4.75)
    top = Inches(3.6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "5000 images par classe. "
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (6.8)
    top = Inches(3.6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Notre service peut vous"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(3.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "fournir une telle base en 7h de tavail comprenant :"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(4.3)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "- 4 h de labellisation assisté par nos algorithmes"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(4.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "- 1 h de process de data augmentation pour vous fournir une base prête à l'emploi"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(5.3)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "- 2 h d'analyse et de rédaction de rapport sur vos données et conseils adaptés à vos compétences en interne"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (0.5)
    top = Inches(5.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Nous vous permettrons donc d'"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (3.22)
    top = Inches(5.8)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "accélérer la phase de recherche et développement de vos projets"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (0.5)
    top = Inches(6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "d'intelligence artificielle"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    left = Inches (2.58)
    top = Inches(6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "en vous donnat toutes les clefs pour la"
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 50, 50)
    
    left = Inches (5.94)
    top = Inches(6)
    width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "réussite de votre projet."
    p.font.size = Pt(14)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(50, 171, 104)
    
    
    
    prs.save('test4.pptx')


